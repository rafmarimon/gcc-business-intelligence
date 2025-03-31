#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
External File Ingestion and Parsing Tool

This script ingests external files (CSV, JSON, DOCX, PDF, etc.) into the system,
parses their content, and associates them with clients for use in reports.
"""

import os
import sys
import json
import argparse
import logging
import time
from datetime import datetime
import hashlib
import shutil
from typing import Dict, List, Any, Optional
import re
import mimetypes

# Ensure proper imports
sys.path.insert(0, os.path.abspath(os.path.dirname(__file__)))

# Import required modules
try:
    from src.models.client_model import ClientModel
    from src.utils.redis_cache import RedisCache
    from dotenv import load_dotenv
except ImportError as e:
    print(f"Error importing required modules: {str(e)}")
    print("Make sure you're running from the project root.")
    sys.exit(1)

# Try to import document parsing libraries
try:
    import pandas as pd
    import PyPDF2
    import docx
    import openpyxl
    import pptx
    PARSERS_AVAILABLE = True
except ImportError:
    print("Warning: Some document parsing libraries are not available.")
    print("Install them with: pip install pandas PyPDF2 python-docx openpyxl python-pptx")
    PARSERS_AVAILABLE = False

# Load environment variables
load_dotenv()

# Configure logging
log_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'logs')
os.makedirs(log_dir, exist_ok=True)

log_filename = os.path.join(log_dir, f'file_ingestion_{datetime.now().strftime("%Y%m%d_%H%M%S")}.log')
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_filename),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class FileIngestor:
    """Handles ingestion and parsing of external files."""
    
    def __init__(self):
        """Initialize the file ingestor."""
        # Initialize Redis cache
        self.redis = RedisCache()
        
        # Initialize client model
        self.client_model = ClientModel()
        
        # Set up storage directory
        data_dir = os.environ.get('DATA_DIR', 'data')
        self.storage_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), data_dir, 'external_files')
        os.makedirs(self.storage_dir, exist_ok=True)
        
        # Define supported file types
        self.supported_types = {
            # Structured data
            '.csv': self._parse_csv,
            '.json': self._parse_json,
            '.xlsx': self._parse_excel,
            '.xls': self._parse_excel,
            
            # Documents
            '.pdf': self._parse_pdf,
            '.docx': self._parse_docx,
            '.doc': self._parse_docx,
            '.txt': self._parse_text,
            '.md': self._parse_text,
            
            # Presentations
            '.pptx': self._parse_pptx,
            '.ppt': self._parse_pptx
        }
        
        logger.info("File ingestor initialized")
    
    def _compute_file_hash(self, file_path: str) -> str:
        """Compute MD5 hash of a file."""
        md5_hash = hashlib.md5()
        with open(file_path, "rb") as f:
            # Read in chunks to handle large files
            for chunk in iter(lambda: f.read(4096), b""):
                md5_hash.update(chunk)
        return md5_hash.hexdigest()
    
    def _detect_file_type(self, file_path: str) -> str:
        """Detect file type based on extension."""
        _, ext = os.path.splitext(file_path)
        return ext.lower()
    
    def _copy_to_storage(self, file_path: str, client_id: str) -> str:
        """Copy file to storage directory and return new path."""
        # Create directory for client if it doesn't exist
        client_dir = os.path.join(self.storage_dir, client_id)
        os.makedirs(client_dir, exist_ok=True)
        
        # Get file hash and extension
        file_hash = self._compute_file_hash(file_path)
        _, ext = os.path.splitext(file_path)
        
        # Create new filename with timestamp and hash
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        new_filename = f"{timestamp}_{file_hash[:8]}{ext}"
        
        # Create destination path
        dest_path = os.path.join(client_dir, new_filename)
        
        # Copy file to storage
        shutil.copy2(file_path, dest_path)
        
        return dest_path
    
    def _parse_csv(self, file_path: str) -> Dict[str, Any]:
        """Parse CSV file."""
        try:
            df = pd.read_csv(file_path)
            
            # Extract basic metadata
            row_count = len(df)
            col_count = len(df.columns)
            
            # Convert to records for storage
            records = df.to_dict('records')
            
            # Create brief preview 
            preview = records[:min(5, row_count)]
            
            return {
                'type': 'csv',
                'row_count': row_count,
                'column_count': col_count,
                'columns': list(df.columns),
                'preview': preview,
                'content': records,
                'summary': f"CSV file with {row_count} rows and {col_count} columns"
            }
            
        except Exception as e:
            logger.error(f"Error parsing CSV file: {str(e)}")
            return {
                'type': 'csv',
                'error': str(e),
                'summary': "Failed to parse CSV file"
            }
    
    def _parse_json(self, file_path: str) -> Dict[str, Any]:
        """Parse JSON file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Determine if it's an array or object
            data_type = 'array' if isinstance(data, list) else 'object'
            
            # Create preview
            if data_type == 'array':
                preview = data[:min(3, len(data))]
                summary = f"JSON array with {len(data)} items"
            else:
                preview = {k: data[k] for k in list(data.keys())[:min(5, len(data.keys()))]}
                summary = f"JSON object with {len(data.keys())} keys"
            
            return {
                'type': 'json',
                'data_type': data_type,
                'preview': preview,
                'content': data,
                'summary': summary
            }
            
        except Exception as e:
            logger.error(f"Error parsing JSON file: {str(e)}")
            return {
                'type': 'json',
                'error': str(e),
                'summary': "Failed to parse JSON file"
            }
    
    def _parse_excel(self, file_path: str) -> Dict[str, Any]:
        """Parse Excel file."""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            sheet_names = excel_file.sheet_names
            
            sheets_data = {}
            total_rows = 0
            
            # Process each sheet
            for sheet in sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet)
                row_count = len(df)
                total_rows += row_count
                
                # Store basic metadata and preview
                sheets_data[sheet] = {
                    'row_count': row_count,
                    'column_count': len(df.columns),
                    'columns': list(df.columns),
                    'preview': df.head(3).to_dict('records')
                }
            
            return {
                'type': 'excel',
                'sheet_count': len(sheet_names),
                'sheets': sheet_names,
                'sheets_data': sheets_data,
                'total_rows': total_rows,
                'summary': f"Excel file with {len(sheet_names)} sheets and {total_rows} total rows"
            }
            
        except Exception as e:
            logger.error(f"Error parsing Excel file: {str(e)}")
            return {
                'type': 'excel',
                'error': str(e),
                'summary': "Failed to parse Excel file"
            }
    
    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """Parse PDF file."""
        try:
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                num_pages = len(pdf_reader.pages)
                
                # Extract text from first few pages for preview
                preview_pages = min(3, num_pages)
                preview_text = []
                
                for i in range(preview_pages):
                    page = pdf_reader.pages[i]
                    text = page.extract_text()
                    preview_text.append({
                        'page': i + 1,
                        'text': text[:500] + ('...' if len(text) > 500 else '')
                    })
                
                # Extract all text for storage
                all_text = []
                for i in range(num_pages):
                    page = pdf_reader.pages[i]
                    text = page.extract_text()
                    all_text.append({
                        'page': i + 1,
                        'text': text
                    })
                
                # Extract metadata if available
                metadata = {}
                if pdf_reader.metadata:
                    for key, value in pdf_reader.metadata.items():
                        if key.startswith('/'):
                            metadata[key[1:]] = value
                
                return {
                    'type': 'pdf',
                    'page_count': num_pages,
                    'preview': preview_text,
                    'content': all_text,
                    'metadata': metadata,
                    'summary': f"PDF document with {num_pages} pages"
                }
                
        except Exception as e:
            logger.error(f"Error parsing PDF file: {str(e)}")
            return {
                'type': 'pdf',
                'error': str(e),
                'summary': "Failed to parse PDF file"
            }
    
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """Parse DOCX file."""
        try:
            doc = docx.Document(file_path)
            
            # Extract paragraphs
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            
            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
            
            return {
                'type': 'docx',
                'paragraph_count': len(paragraphs),
                'table_count': len(tables),
                'preview': ' '.join(paragraphs[:3])[:500] + ('...' if len(' '.join(paragraphs[:3])) > 500 else ''),
                'content': {
                    'paragraphs': paragraphs,
                    'tables': tables
                },
                'summary': f"Word document with {len(paragraphs)} paragraphs and {len(tables)} tables"
            }
            
        except Exception as e:
            logger.error(f"Error parsing DOCX file: {str(e)}")
            return {
                'type': 'docx',
                'error': str(e),
                'summary': "Failed to parse Word document"
            }
    
    def _parse_text(self, file_path: str) -> Dict[str, Any]:
        """Parse plain text or markdown file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Split into lines
            lines = content.split('\n')
            
            # Determine if it's markdown by checking for common markdown features
            is_markdown = any(line.startswith('#') for line in lines) or \
                          '**' in content or \
                          '*' in content or \
                          '[' in content and '](' in content
            
            file_type = 'markdown' if is_markdown else 'text'
            
            return {
                'type': file_type,
                'line_count': len(lines),
                'char_count': len(content),
                'preview': content[:500] + ('...' if len(content) > 500 else ''),
                'content': content,
                'summary': f"{'Markdown' if is_markdown else 'Text'} file with {len(lines)} lines and {len(content)} characters"
            }
            
        except Exception as e:
            logger.error(f"Error parsing text file: {str(e)}")
            return {
                'type': 'text',
                'error': str(e),
                'summary': "Failed to parse text file"
            }
    
    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """Parse PowerPoint presentation."""
        try:
            presentation = pptx.Presentation(file_path)
            
            slides_data = []
            for i, slide in enumerate(presentation.slides):
                slide_data = {
                    'slide_number': i + 1,
                    'shapes': len(slide.shapes),
                    'text': []
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        slide_data['text'].append(shape.text)
                
                slides_data.append(slide_data)
            
            return {
                'type': 'pptx',
                'slide_count': len(presentation.slides),
                'preview': slides_data[:3],
                'content': slides_data,
                'summary': f"PowerPoint presentation with {len(presentation.slides)} slides"
            }
            
        except Exception as e:
            logger.error(f"Error parsing PowerPoint file: {str(e)}")
            return {
                'type': 'pptx',
                'error': str(e),
                'summary': "Failed to parse PowerPoint presentation"
            }
    
    def _extract_keywords(self, content: Dict[str, Any]) -> List[str]:
        """Extract keywords from content."""
        keywords = []
        
        # Handle text-based content
        text = ""
        content_type = content.get('type', '')
        
        if content_type == 'csv':
            # For CSV, concatenate column names and first few rows
            columns = content.get('columns', [])
            preview = content.get('preview', [])
            text = ' '.join(columns)
            for row in preview:
                text += ' ' + ' '.join(str(val) for val in row.values())
                
        elif content_type == 'json':
            # For JSON, convert to string
            preview = content.get('preview', {})
            text = json.dumps(preview)
            
        elif content_type in ['excel', 'pdf', 'docx', 'text', 'markdown', 'pptx']:
            # For other types, use preview or content
            preview = content.get('preview', '')
            if isinstance(preview, str):
                text = preview
            elif isinstance(preview, list):
                if content_type == 'pdf':
                    text = ' '.join(page.get('text', '') for page in preview)
                elif content_type == 'pptx':
                    for slide in preview:
                        text += ' '.join(slide.get('text', []))
                else:
                    text = ' '.join(str(item) for item in preview)
            elif isinstance(preview, dict):
                text = json.dumps(preview)
        
        # Simple keyword extraction using frequency
        if text:
            # Convert to lowercase and remove punctuation
            text = re.sub(r'[^\w\s]', ' ', text.lower())
            
            # Split into words and count frequency
            words = text.split()
            word_freq = {}
            for word in words:
                if len(word) > 3:  # Ignore short words
                    word_freq[word] = word_freq.get(word, 0) + 1
            
            # Get top keywords by frequency
            sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
            keywords = [word for word, freq in sorted_words[:20]]  # Top 20 keywords
        
        return keywords
    
    def _categorize_content(self, content: Dict[str, Any], client: Dict[str, Any]) -> List[str]:
        """Categorize content based on client interests."""
        categories = []
        
        # Get client interests
        interests = client.get('interests_list', [])
        
        # Extract keywords from content
        keywords = self._extract_keywords(content)
        
        # Check for overlaps with client interests
        for interest in interests:
            interest_lower = interest.lower()
            
            # Check if interest appears in keywords
            if any(interest_lower in keyword.lower() for keyword in keywords):
                categories.append(interest)
        
        # Add file type as a category
        file_type = content.get('type', 'unknown')
        categories.append(f"file_type:{file_type}")
        
        return categories
    
    def ingest_file(self, file_path: str, client_name_or_id: str, tags: Optional[List[str]] = None) -> Dict[str, Any]:
        """
        Ingest a file into the system.
        
        Args:
            file_path: Path to the file to ingest
            client_name_or_id: Client name or ID to associate the file with
            tags: Optional tags to apply to the file
            
        Returns:
            Dictionary with ingestion results
        """
        # Check if file exists
        if not os.path.isfile(file_path):
            logger.error(f"File not found: {file_path}")
            return {
                'success': False,
                'error': f"File not found: {file_path}"
            }
        
        # Get client
        client = self._get_client(client_name_or_id)
        if not client:
            return {
                'success': False,
                'error': f"Client not found: {client_name_or_id}"
            }
        
        client_id = client.get('id')
        
        # Detect file type
        file_type = self._detect_file_type(file_path)
        if file_type not in self.supported_types:
            logger.warning(f"Unsupported file type: {file_type}")
            return {
                'success': False,
                'error': f"Unsupported file type: {file_type}"
            }
        
        # Copy file to storage
        storage_path = self._copy_to_storage(file_path, client_id)
        logger.info(f"File copied to storage: {storage_path}")
        
        # Parse file content
        parse_method = self.supported_types[file_type]
        content = parse_method(file_path)
        
        # Add file metadata
        file_size = os.path.getsize(file_path)
        file_name = os.path.basename(file_path)
        
        # Create file record
        file_record = {
            'client_id': client_id,
            'client_name': client.get('name'),
            'file_name': file_name,
            'file_path': storage_path,
            'file_type': file_type,
            'file_size': file_size,
            'ingested_at': datetime.now().isoformat(),
            'content': content,
            'tags': tags or []
        }
        
        # Categorize content
        categories = self._categorize_content(content, client)
        file_record['categories'] = categories
        
        # Store in Redis
        file_id = f"external_file:{client_id}:{datetime.now().strftime('%Y%m%d%H%M%S')}"
        self.redis.set(file_id, file_record)
        
        # Associate with client
        client_files = self.redis.get(f"client:{client_id}:external_files") or []
        client_files.append(file_id)
        self.redis.set(f"client:{client_id}:external_files", client_files)
        
        logger.info(f"File ingested successfully: {file_id}")
        
        return {
            'success': True,
            'file_id': file_id,
            'client_id': client_id,
            'client_name': client.get('name'),
            'file_name': file_name,
            'file_type': file_type,
            'file_size': file_size,
            'storage_path': storage_path,
            'categories': categories,
            'tags': tags or [],
            'summary': content.get('summary', 'File processed successfully')
        }
    
    def ingest_directory(self, directory_path: str, client_name_or_id: str, recursive: bool = False) -> Dict[str, Any]:
        """
        Ingest all files in a directory.
        
        Args:
            directory_path: Path to the directory to ingest
            client_name_or_id: Client name or ID to associate the files with
            recursive: Whether to recursively process subdirectories
            
        Returns:
            Dictionary with ingestion results
        """
        # Check if directory exists
        if not os.path.isdir(directory_path):
            logger.error(f"Directory not found: {directory_path}")
            return {
                'success': False,
                'error': f"Directory not found: {directory_path}"
            }
        
        # Get client
        client = self._get_client(client_name_or_id)
        if not client:
            return {
                'success': False,
                'error': f"Client not found: {client_name_or_id}"
            }
        
        # Process files
        results = {
            'success': True,
            'client_id': client.get('id'),
            'client_name': client.get('name'),
            'directory': directory_path,
            'processed_files': [],
            'failed_files': [],
            'skipped_files': []
        }
        
        # Walk through directory
        for root, dirs, files in os.walk(directory_path):
            for file in files:
                file_path = os.path.join(root, file)
                
                # Check if file type is supported
                file_type = self._detect_file_type(file_path)
                if file_type not in self.supported_types:
                    results['skipped_files'].append({
                        'file': file_path,
                        'reason': f"Unsupported file type: {file_type}"
                    })
                    continue
                
                # Ingest file
                ingest_result = self.ingest_file(file_path, client.get('id'))
                
                if ingest_result.get('success'):
                    results['processed_files'].append({
                        'file': file_path,
                        'file_id': ingest_result.get('file_id'),
                        'file_type': file_type,
                        'summary': ingest_result.get('summary')
                    })
                else:
                    results['failed_files'].append({
                        'file': file_path,
                        'error': ingest_result.get('error')
                    })
            
            # Stop if not recursive
            if not recursive:
                break
        
        # Update result counts
        results['total_files'] = len(results['processed_files']) + len(results['failed_files']) + len(results['skipped_files'])
        results['processed_count'] = len(results['processed_files'])
        results['failed_count'] = len(results['failed_files'])
        results['skipped_count'] = len(results['skipped_files'])
        
        logger.info(f"Directory ingestion complete: {results['processed_count']} processed, {results['failed_count']} failed, {results['skipped_count']} skipped")
        
        return results
    
    def _get_client(self, client_name_or_id: str) -> Optional[Dict[str, Any]]:
        """Get client by name or ID."""
        # Try to get client by ID first
        client = self.client_model.get_client_by_id(client_name_or_id)
        
        # If not found, try by name
        if not client:
            clients = self.client_model.get_all_clients()
            for c in clients:
                if c.get('name', '').lower() == client_name_or_id.lower():
                    client = c
                    break
        
        if not client:
            logger.error(f"Client not found: {client_name_or_id}")
            return None
        
        return client
    
    def list_files(self, client_name_or_id: str) -> List[Dict[str, Any]]:
        """List all files associated with a client."""
        client = self._get_client(client_name_or_id)
        if not client:
            return []
        
        client_id = client.get('id')
        
        # Get file IDs
        file_ids = self.redis.get(f"client:{client_id}:external_files") or []
        
        # Get file details
        files = []
        for file_id in file_ids:
            file_data = self.redis.get(file_id)
            if file_data:
                files.append({
                    'file_id': file_id,
                    'file_name': file_data.get('file_name'),
                    'file_type': file_data.get('file_type'),
                    'file_size': file_data.get('file_size'),
                    'ingested_at': file_data.get('ingested_at'),
                    'categories': file_data.get('categories', []),
                    'tags': file_data.get('tags', []),
                    'summary': file_data.get('content', {}).get('summary', 'File processed successfully')
                })
        
        return files
    
    def get_file(self, file_id: str) -> Optional[Dict[str, Any]]:
        """Get file details by ID."""
        file_data = self.redis.get(file_id)
        if not file_data:
            logger.error(f"File not found: {file_id}")
            return None
        
        return file_data
    
    def delete_file(self, file_id: str) -> bool:
        """Delete a file from the system."""
        file_data = self.redis.get(file_id)
        if not file_data:
            logger.error(f"File not found: {file_id}")
            return False
        
        # Get client ID
        client_id = file_data.get('client_id')
        
        # Get storage path
        storage_path = file_data.get('file_path')
        
        # Delete physical file
        if storage_path and os.path.exists(storage_path):
            try:
                os.remove(storage_path)
                logger.info(f"File deleted from storage: {storage_path}")
            except Exception as e:
                logger.error(f"Error deleting file from storage: {str(e)}")
        
        # Remove from client association
        client_files = self.redis.get(f"client:{client_id}:external_files") or []
        if file_id in client_files:
            client_files.remove(file_id)
            self.redis.set(f"client:{client_id}:external_files", client_files)
        
        # Delete from Redis
        self.redis.delete(file_id)
        
        logger.info(f"File deleted: {file_id}")
        
        return True

def main():
    """Main function to parse arguments and execute commands."""
    parser = argparse.ArgumentParser(description='External File Ingestion and Parsing Tool')
    
    # Command selection
    cmd_group = parser.add_mutually_exclusive_group(required=True)
    cmd_group.add_argument('--file', type=str, help='Path to a file to ingest')
    cmd_group.add_argument('--dir', type=str, help='Path to a directory to ingest')
    cmd_group.add_argument('--list', action='store_true', help='List files for a client')
    cmd_group.add_argument('--get', type=str, help='Get file details by ID')
    cmd_group.add_argument('--delete', type=str, help='Delete a file by ID')
    
    # Client information
    parser.add_argument('--client', type=str, help='Client name or ID')
    
    # Tags
    parser.add_argument('--tags', type=str, help='Comma-separated list of tags to apply to the file')
    
    # Directory options
    parser.add_argument('--recursive', action='store_true', help='Recursively process subdirectories')
    
    # Output options
    parser.add_argument('--output', type=str, help='Output file for results (JSON format)')
    
    args = parser.parse_args()
    
    # Initialize ingestor
    ingestor = FileIngestor()
    
    # Execute command
    result = None
    
    # Ingest file
    if args.file:
        if not args.client:
            print("Error: Client is required when ingesting a file.")
            return 1
        
        # Parse tags
        tags = args.tags.split(',') if args.tags else None
        
        # Ingest file
        result = ingestor.ingest_file(args.file, args.client, tags)
        
        if result.get('success'):
            print(f"File ingested successfully: {result.get('file_id')}")
            print(f"File type: {result.get('file_type')}")
            print(f"Summary: {result.get('summary')}")
        else:
            print(f"Error ingesting file: {result.get('error')}")
            return 1
    
    # Ingest directory
    elif args.dir:
        if not args.client:
            print("Error: Client is required when ingesting a directory.")
            return 1
        
        # Ingest directory
        result = ingestor.ingest_directory(args.dir, args.client, args.recursive)
        
        if result.get('success'):
            print(f"Directory ingestion complete:")
            print(f"  Processed: {result.get('processed_count')}")
            print(f"  Failed: {result.get('failed_count')}")
            print(f"  Skipped: {result.get('skipped_count')}")
        else:
            print(f"Error ingesting directory: {result.get('error')}")
            return 1
    
    # List files
    elif args.list:
        if not args.client:
            print("Error: Client is required when listing files.")
            return 1
        
        # List files
        files = ingestor.list_files(args.client)
        
        if files:
            print(f"Files for client {args.client}:")
            for file in files:
                print(f"  {file.get('file_name')} ({file.get('file_type')}) - {file.get('file_id')}")
                print(f"    Ingested: {file.get('ingested_at')}")
                print(f"    Summary: {file.get('summary')}")
                print(f"    Categories: {', '.join(file.get('categories', []))}")
                print(f"    Tags: {', '.join(file.get('tags', []))}")
                print()
        else:
            print(f"No files found for client {args.client}")
        
        result = {'files': files}
    
    # Get file details
    elif args.get:
        # Get file
        file = ingestor.get_file(args.get)
        
        if file:
            print(f"File details for {args.get}:")
            print(f"  Name: {file.get('file_name')}")
            print(f"  Type: {file.get('file_type')}")
            print(f"  Size: {file.get('file_size')} bytes")
            print(f"  Ingested: {file.get('ingested_at')}")
            print(f"  Client: {file.get('client_name')} ({file.get('client_id')})")
            print(f"  Path: {file.get('file_path')}")
            print(f"  Categories: {', '.join(file.get('categories', []))}")
            print(f"  Tags: {', '.join(file.get('tags', []))}")
            
            content = file.get('content', {})
            print(f"  Summary: {content.get('summary', 'No summary available')}")
        else:
            print(f"File not found: {args.get}")
            return 1
        
        result = file
    
    # Delete file
    elif args.delete:
        # Delete file
        success = ingestor.delete_file(args.delete)
        
        if success:
            print(f"File deleted: {args.delete}")
        else:
            print(f"Error deleting file: {args.delete}")
            return 1
        
        result = {'success': success, 'file_id': args.delete}
    
    # Save output to file if requested
    if args.output and result:
        try:
            with open(args.output, 'w', encoding='utf-8') as f:
                json.dump(result, f, indent=2)
            print(f"Results saved to {args.output}")
        except Exception as e:
            print(f"Error saving results to file: {str(e)}")
    
    return 0

if __name__ == '__main__':
    sys.exit(main()) 